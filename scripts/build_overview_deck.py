"""
Build memory-bank-overview.pptx with one consistent visual system that
mirrors training/presentation.html — dark blue background, magenta brand
rail, magenta title underline, magenta T·MOBILE top-right, white body text.

Runs from repo root:  python scripts/build_overview_deck.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

REPO = Path(__file__).resolve().parent.parent
OUTPUT = REPO / "memory-bank-overview.pptx"
T_ICON = REPO / "brand" / "tmobile-t-icon.png"          # small: corner watermark
T_ICON_LARGE = REPO / "brand" / "tmobile-t-icon-large.png"  # large: title-slide hero

TEAMS_CHANNEL_NAME = "RE - SkyNet Support - AI Discussion"
TEAMS_CHANNEL_URL = (
    "https://teams.microsoft.com/l/channel/"
    "19%3A7130c6f6eb354efda1d4b3fa89546215%40thread.tacv2/"
    "RE%20-%20SkyNet%20Support%20-%20AI%20Discussion"
    "?groupId=4f72c46d-e46e-43b9-a3d6-1de811294cf8"
    "&tenantId=be0f980b-dd99-4b19-bd7b-bc71a09b026c"
)

# Set by main() once BUILDERS length is known; read by chrome footer.
TOTAL_SLIDES = 0

# Auto-incremented by add_blank_slide; reset by main().
_slide_counter = [0]

# T-Mobile brand palette (from T-Mobile Branding MCP)
MAGENTA = RGBColor(0xE2, 0x00, 0x74)
BERRY = RGBColor(0x86, 0x1B, 0x54)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x6A, 0x6A, 0x6A)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)

# Web-deck palette (matches training/presentation.html)
BG_NAVY = RGBColor(0x1A, 0x1A, 0x2E)
BG_NAVY_ALT = RGBColor(0x16, 0x21, 0x3E)
CODE_BG = RGBColor(0x0D, 0x11, 0x17)
BODY_TEXT = WHITE
SUB_TEXT = LIGHT_GRAY

FONT = "TeleNeo"
FONT_MONO = "Consolas"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

FOOTER_Y = 7.15
MAX_CONTENT_Y = 7.10


def _check_bounds(top_in, height_in, kind):
    # Fail loudly if body content would overrun the footer band.
    # top_in > 1.5 exempts the title/subtitle zone at y~0.3.
    # 0.005" epsilon absorbs IEEE-754 addition noise (e.g., 6.70 + 0.40 = 7.1000000000000005).
    if top_in > 1.5 and top_in + height_in > MAX_CONTENT_Y + 0.005:
        raise ValueError(
            f"{kind} at y={top_in:.2f}, h={height_in:.2f} ends at "
            f"{top_in + height_in:.4f} - exceeds MAX_CONTENT_Y={MAX_CONTENT_Y}. "
            f"Fix the slide layout instead of raising this bound."
        )


def set_slide_background(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_slide_chrome(slide, slide_num: int, *, rail_width_in: float = 0.1,
                     footer_left_in: float = 0.25):
    """Dark navy bg + magenta left rail + T icon top-right + slide-num footer.

    `rail_width_in` — width of the magenta left accent (Slide 1 uses a hero bar).
    `footer_left_in` — where the footer text starts; must clear the rail on wide-rail slides.
    """
    set_slide_background(slide, BG_NAVY if slide_num % 2 == 1 else BG_NAVY_ALT)

    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                  Inches(rail_width_in), Inches(SLIDE_H_IN))
    rail.line.fill.background()
    rail.fill.solid()
    rail.fill.fore_color.rgb = MAGENTA

    # T-Mobile "T" icon in the top-right (extracted from the official brand template)
    if T_ICON.exists():
        icon_w = 0.5
        slide.shapes.add_picture(
            str(T_ICON),
            Inches(SLIDE_W_IN - icon_w - 0.35),
            Inches(0.28),
            Inches(icon_w),
            Inches(icon_w),
        )

    # Footer: "Memory Bank Standard | Internal Enablement   ·   N / 16"
    footer = slide.shapes.add_textbox(
        Inches(footer_left_in), Inches(SLIDE_H_IN - 0.35),
        Inches(SLIDE_W_IN - footer_left_in - 0.25), Inches(0.25)
    )
    tf = footer.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = f"Memory Bank Standard  ·  Internal Enablement  ·  {slide_num} / {TOTAL_SLIDES}"
    r.font.name = FONT
    r.font.size = Pt(9)
    r.font.color.rgb = DARK_GRAY


def add_blank_slide(prs: Presentation, slide_num: int = None, *,
                    rail_width_in: float = 0.1, footer_left_in: float = 0.25):
    if slide_num is None:
        _slide_counter[0] += 1
        slide_num = _slide_counter[0]
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # true blank layout
    add_slide_chrome(slide, slide_num,
                     rail_width_in=rail_width_in, footer_left_in=footer_left_in)
    return slide


def set_run(run, text, *, size=18, bold=False, color=BODY_TEXT, font=FONT):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(
    slide,
    left_in,
    top_in,
    width_in,
    height_in,
    text,
    *,
    size=18,
    bold=False,
    color=BODY_TEXT,
    font=FONT,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    _check_bounds(top_in, height_in, "add_text")
    box = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(36000)
    tf.margin_bottom = Emu(36000)
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size=size, bold=bold, color=color, font=font)
    return box


def add_bullets(
    slide,
    left_in,
    top_in,
    width_in,
    height_in,
    items,
    *,
    size=16,
    color=BODY_TEXT,
    bullet_char="•",
    line_spacing=1.2,
    space_after_pt=4,
):
    _check_bounds(top_in, height_in, "add_bullets")
    box = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after_pt)
        parts = item if isinstance(item, list) else [(item, {})]
        first_opts = parts[0][1] if parts else {}
        no_bullet = first_opts.get("no_bullet", False)
        if not no_bullet:
            set_run(p.add_run(), f"{bullet_char}  ",
                    size=size, color=MAGENTA, bold=True)
        for text, opts in parts:
            run = p.add_run()
            set_run(
                run,
                text,
                size=opts.get("size", size),
                bold=opts.get("bold", False),
                color=opts.get("color", color),
                font=opts.get("font", FONT),
            )
    return box


def add_title(slide, title, *, subtitle=None):
    """Magenta title with magenta underline — matches HTML h2 treatment."""
    # Title text
    add_text(slide, 0.55, 0.55, SLIDE_W_IN - 1.1, 0.85, title,
             size=32, bold=True, color=MAGENTA)
    # Magenta underline (thin bar) right under the title text
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.3), Inches(4.2), Inches(0.05)
    )
    underline.line.fill.background()
    underline.fill.solid()
    underline.fill.fore_color.rgb = MAGENTA
    if subtitle:
        add_text(slide, 0.55, 1.45, SLIDE_W_IN - 1.1, 0.5, subtitle,
                 size=16, color=SUB_TEXT)


def add_rounded_box(
    slide, left_in, top_in, width_in, height_in, *,
    fill=CODE_BG, line=None, accent_top=False, accent_left=False,
):
    _check_bounds(top_in, height_in, "add_rounded_box")
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in),
    )
    shape.adjustments[0] = 0.05
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    shape.shadow.inherit = False

    if accent_top:
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_in), Inches(top_in), Inches(width_in), Inches(0.05),
        )
        top_bar.line.fill.background()
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = MAGENTA
    if accent_left:
        left_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_in), Inches(top_in), Inches(0.05), Inches(height_in),
        )
        left_bar.line.fill.background()
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = MAGENTA
    return shape


def add_code_box(slide, left_in, top_in, width_in, height_in, code, *, size=16):
    _check_bounds(top_in, height_in, "add_code_box")
    box = add_rounded_box(slide, left_in, top_in, width_in, height_in,
                          fill=CODE_BG, accent_left=True)
    tb = slide.shapes.add_textbox(
        Inches(left_in + 0.18), Inches(top_in + 0.12),
        Inches(width_in - 0.3), Inches(height_in - 0.24),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(18000)
    tf.margin_right = Emu(18000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    set_run(p.add_run(), code, size=size, bold=True, color=WHITE, font=FONT_MONO)
    return box


def set_notes(slide, talk_track, files_referenced=None, assumptions=None, missing=None):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    p0 = notes.paragraphs[0]
    p0.text = "Talk track:"
    p0.runs[0].font.bold = True
    p0.runs[0].font.size = Pt(12)
    for line in talk_track.strip().split("\n"):
        p = notes.add_paragraph()
        p.text = line.strip()
        p.runs[0].font.size = Pt(11)

    for label, items in (("Repo files referenced:", files_referenced),
                         ("Assumptions:", assumptions),
                         ("Missing info / follow-ups:", missing)):
        if not items:
            continue
        p = notes.add_paragraph()
        p.text = ""
        p = notes.add_paragraph()
        p.text = label
        p.runs[0].font.bold = True
        p.runs[0].font.size = Pt(12)
        for x in items:
            p = notes.add_paragraph()
            p.text = f"- {x}"
            p.runs[0].font.size = Pt(11)


# =============================================================================
# Slide builders — 16 slides, one consistent look
# =============================================================================


def build_slide_1_title(prs):
    # Standard chrome — rail and footer match every other slide
    slide = add_blank_slide(prs)

    # Faint magenta circle behind the brain for extra hero weight
    halo = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(SLIDE_W_IN / 2 - 1.4), Inches(0.95),
        Inches(2.8), Inches(2.8),
    )
    halo.line.fill.background()
    halo.fill.solid()
    halo.fill.fore_color.rgb = BERRY   # deep magenta — sits behind the emoji

    # Brain emoji hero — big, centered
    brain = slide.shapes.add_textbox(
        Inches(SLIDE_W_IN / 2 - 1.5), Inches(1.15),
        Inches(3.0), Inches(2.5),
    )
    tf = brain.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "🧠"
    r.font.name = "Segoe UI Emoji"
    r.font.size = Pt(150)

    # Magenta divider bar centered under the brain
    bar_w = 2.2
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(SLIDE_W_IN / 2 - bar_w / 2), Inches(4.05),
        Inches(bar_w), Inches(0.07),
    )
    divider.line.fill.background()
    divider.fill.solid()
    divider.fill.fore_color.rgb = MAGENTA

    # Title, subtitle, tagline — all centered
    add_text(slide, 0.6, 4.3, SLIDE_W_IN - 1.2, 1.1,
             "Memory-Bank Project",
             size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, 0.6, 5.4, SLIDE_W_IN - 1.2, 0.7,
             "Safer, more consistent AI coding assistants — across sessions.",
             size=22, color=SUB_TEXT, align=PP_ALIGN.CENTER)
    add_text(slide, 0.6, 6.15, SLIDE_W_IN - 1.2, 0.45,
             "Internal enablement  ·  T-Mobile Release Engineering / AERO",
             size=13, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    set_notes(
        slide,
        """Onboarding walkthrough of the Memory-Bank project. Open with the problem: AI assistants lose context between sessions. Memory-Bank is a template repository that fixes that — productive in under 15 minutes per project.""",
        files_referenced=["README.md line 3"],
    )


def build_slide_2_audience(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Who this is for")
    items = [
        [("Individual contributors", {"bold": True, "color": MAGENTA}),
         ("  — engineers who use Cursor or Claude Code daily.", {})],
        [("Tech leads", {"bold": True, "color": MAGENTA}),
         ("  — owners of AI coding standards; reviewers who want predictable output.", {})],
        [("Engineering managers", {"bold": True, "color": MAGENTA}),
         ("  — stakeholders who need safer defaults, fewer incidents, faster onboarding.", {})],
        [("Anyone AI-assisting a repo", {"bold": True, "color": MAGENTA}),
         ("  — mixed skill levels welcome; install slides are plain-English.", {})],
    ]
    add_bullets(slide, 0.7, 2.2, SLIDE_W_IN - 1.4, 3.8, items, size=18)
    add_text(slide, 0.7, 6.1, SLIDE_W_IN - 1.4, 0.5,
             "If you touch code in Cursor or Claude Code, this deck is for you.",
             size=15, color=MAGENTA, bold=True)
    set_notes(
        slide,
        """Set audience expectations. Nobody has to be a deep technical expert — install slides are written for that.""",
        files_referenced=["README.md lines 138–148 (IDE support)"],
    )


def build_slide_3_exec_summary(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Executive summary")
    items = [
        [("What it is", {"bold": True, "color": MAGENTA}),
         ("  — a template repository you drop into any project.", {})],
        [("What it does", {"bold": True, "color": MAGENTA}),
         ("  — makes Cursor and Claude Code read persistent project context, follow a 3-tier security model, enforce code-quality and logging standards, run a structured feature workflow, and orchestrate a multi-agent code review.", {})],
        [("Why it matters", {"bold": True, "color": MAGENTA}),
         ("  — fewer secrets committed, fewer inconsistent rewrites, fewer missed tests, faster onboarding, and review coverage that doesn't depend on one reviewer catching everything.", {})],
        [("How fast", {"bold": True, "color": MAGENTA}),
         ("  — productive in under 15 minutes per project. One-liner install on Windows or Mac.", {})],
    ]
    add_bullets(slide, 0.7, 2.1, SLIDE_W_IN - 1.4, 3.9, items, size=17)
    add_rounded_box(slide, 0.7, 6.05, SLIDE_W_IN - 1.4, 0.85,
                    fill=CODE_BG, accent_left=True)
    add_text(slide, 1.0, 6.15, SLIDE_W_IN - 2.0, 0.7,
             "Five standards · IDE rule files · setup scripts · training materials · multi-agent /code-review.",
             size=15, bold=True, color=WHITE)
    set_notes(
        slide,
        """One-minute pitch. Anchor on 15 minutes per project.""",
        files_referenced=["README.md lines 1–3, 32–78"],
    )


def build_slide_4_problem(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "The problem",
              subtitle="AI assistants forget. Standards drift. One-pass reviews miss things.")

    # Left panel — the pain
    add_rounded_box(slide, 0.7, 2.25, 6.0, 4.3, fill=CODE_BG, accent_top=True)
    add_text(slide, 0.9, 2.42, 5.6, 0.5, "Without memory-bank",
             size=18, bold=True, color=MAGENTA)
    items_bad = [
        [("AI forgets context", {"bold": True, "color": WHITE}),
         ("  between sessions; you re-explain the stack every time.", {"color": SUB_TEXT})],
        [("Different results per user", {"bold": True, "color": WHITE}),
         ("  — one dev gets secure code, another commits a secret.", {"color": SUB_TEXT})],
        [("Standards drift", {"bold": True, "color": WHITE}),
         ("  — security, logging, quality inconsistently applied.", {"color": SUB_TEXT})],
        [("One-pass reviews miss things", {"bold": True, "color": WHITE}),
         ("  — a single reviewer can't cover everything equally.", {"color": SUB_TEXT})],
    ]
    add_bullets(slide, 0.9, 3.0, 5.6, 3.4, items_bad, size=13, bullet_char="✕")

    # Right panel — the fix
    add_rounded_box(slide, 6.9, 2.25, 6.0, 4.3, fill=CODE_BG, accent_top=True)
    add_text(slide, 7.1, 2.42, 5.6, 0.5, "With memory-bank",
             size=18, bold=True, color=MAGENTA)
    items_good = [
        [("Context persists", {"bold": True, "color": WHITE}),
         ("  — 5 files the AI reads at every session start.", {"color": SUB_TEXT})],
        [("Same defaults for everyone", {"bold": True, "color": WHITE}),
         ("  — rule files install globally or per-project.", {"color": SUB_TEXT})],
        [("Standards auto-enforced", {"bold": True, "color": WHITE}),
         ("  — security, quality, logging, workflow.", {"color": SUB_TEXT})],
        [("Multi-agent /code-review", {"bold": True, "color": WHITE}),
         ("  — parallel role reviewers + tests + auditor compare.", {"color": SUB_TEXT})],
    ]
    add_bullets(slide, 7.1, 3.0, 5.6, 3.4, items_good, size=13)

    set_notes(
        slide,
        """Concrete before/after. Tease the multi-agent code review — slides 11 and 12 go deep on it.""",
        files_referenced=["README.md lines 88–104", "templates/CLAUDE.md Context Compaction"],
    )


def build_slide_5_solution(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Solution overview", subtitle="Nine pieces, one template repository.")
    pieces = [
        ("1", "Persistent memory-bank", "Five structured files the AI reads first."),
        ("2", "3-tier security guardrails", "BLOCK / CONFIRM / WARN classifications."),
        ("3", "Code quality standards", "Run tests before done; WHY-only comments."),
        ("4", "Accessibility", "WCAG 2.1 AA — glob-scoped Cursor rule + /accessibility-review audit."),
        ("5", "Structured logging", "JSON, correlation IDs, auto-PII redaction."),
        ("6", "Feature workflow", "Brainstorm → spec → plan → implement → review."),
        ("7", "IDE rule files", "Cursor .mdc + Claude Code CLAUDE.md / slash commands."),
        ("8", "Setup scripts", "Windows/Mac/Linux one-liners."),
        ("9", "Multi-agent /code-review", "Security · Performance · Style + tests + auditor."),
    ]
    col_w = 4.05
    row_h = 1.2
    left0 = 0.65
    top0 = 2.3
    for i, (num, title, desc) in enumerate(pieces):
        col = i % 3
        row = i // 3
        left = left0 + col * (col_w + 0.08)
        top = top0 + row * (row_h + 0.12)
        add_rounded_box(slide, left, top, col_w, row_h, fill=CODE_BG, accent_left=True)
        # number circle
        add_text(slide, left + 0.22, top + 0.14, 0.45, 0.4, num,
                 size=18, bold=True, color=MAGENTA)
        add_text(slide, left + 0.75, top + 0.12, col_w - 0.9, 0.4, title,
                 size=14, bold=True, color=WHITE)
        add_text(slide, left + 0.75, top + 0.52, col_w - 0.9, 0.6, desc,
                 size=11, color=SUB_TEXT)
    set_notes(
        slide,
        """Don't walk every tile. Point at the grid: 'nine pieces, one repo.'""",
        files_referenced=["README.md lines 32–78"],
    )


def build_slide_6_memory_bank(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Persistent memory-bank",
              subtitle="A project notebook the AI reads before helping.")
    files = [
        ("projectbrief.md", "Non-negotiable requirements and constraints.", "Rarely"),
        ("systemPatterns.md", "Architecture decisions and patterns to follow.", "New patterns"),
        ("techContext.md", "Tech stack, dependencies, environment.", "Tech changes"),
        ("activeContext.md", "Current focus and immediate next steps.", "Every session"),
        ("progress.md", "What's complete, in progress, planned.", "Milestones"),
    ]
    col_x = [0.7, 4.1, 10.3]
    col_w = [3.3, 6.1, 2.4]
    top = 2.3
    row_h = 0.6
    # Header row (magenta band)
    header = add_rounded_box(slide, col_x[0], top, sum(col_w) + 0.1, row_h - 0.05,
                             fill=MAGENTA)
    add_text(slide, col_x[0] + 0.15, top + 0.08, col_w[0], 0.4, "File",
             size=14, bold=True, color=WHITE)
    add_text(slide, col_x[1] + 0.1, top + 0.08, col_w[1], 0.4, "Purpose",
             size=14, bold=True, color=WHITE)
    add_text(slide, col_x[2] + 0.1, top + 0.08, col_w[2], 0.4, "Update when",
             size=14, bold=True, color=WHITE)
    for i, (name, purpose, when) in enumerate(files):
        top_r = top + row_h + i * row_h
        if i % 2 == 0:
            add_rounded_box(slide, col_x[0], top_r, sum(col_w) + 0.1, row_h - 0.05,
                            fill=CODE_BG)
        add_text(slide, col_x[0] + 0.15, top_r + 0.1, col_w[0], 0.4, name,
                 size=13, bold=True, color=MAGENTA, font=FONT_MONO)
        add_text(slide, col_x[1] + 0.1, top_r + 0.1, col_w[1], 0.4, purpose,
                 size=13, color=WHITE)
        add_text(slide, col_x[2] + 0.1, top_r + 0.1, col_w[2], 0.4, when,
                 size=13, color=SUB_TEXT)

    add_rounded_box(slide, 0.7, 5.95, SLIDE_W_IN - 1.4, 0.70,
                    fill=CODE_BG, accent_left=True)
    add_text(slide, 1.0, 6.02, SLIDE_W_IN - 2.0, 0.55,
             'Rule: "At the start of every conversation (and after context compaction), silently read ALL files in memory-bank/." (templates/CLAUDE.md, line 7)',
             size=13, color=WHITE)
    set_notes(
        slide,
        """The AI 'reads first' effect is what makes it feel like it already knows your project. The next slide covers Handoff — how context is preserved when sessions fill up.""",
        files_referenced=["templates/memory-bank/*.md", "templates/CLAUDE.md lines 5–20"],
    )


def build_slide_7_security(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "3-tier security guardrails",
              subtitle="Stop dangerous defaults before they ship.")
    tiers = [
        ("BLOCK", "I will refuse", MAGENTA, [
            "Commit secrets (.env, API keys, *.pem, *.key)",
            "Force-push to main/master",
            "Destructive system commands",
            "Slopsquat packages (unverified on registry)",
            "Hardcode credentials in mcp.json",
            "Expose secrets in logs",
        ]),
        ("CONFIRM", "I will ask first", BERRY, [
            "Delete any files",
            "Bulk operations on > 3 files",
            "Amend commits / skip git hooks",
            "Force-push any branch",
            "DROP / DELETE without WHERE / TRUNCATE",
            "Modify CI/CD or auth/permission files",
        ]),
        ("WARN", "I will note the risk", DARK_GRAY, [
            "Large changes (> 5 files or > 200 lines)",
            "Creating new files when editing works",
            "Missing tests for new functionality",
            "Skipping verification before completion",
        ]),
    ]
    col_w = 4.05
    top = 2.2
    for i, (name, tag, banner, items) in enumerate(tiers):
        left = 0.65 + i * (col_w + 0.08)
        add_rounded_box(slide, left, top, col_w, 4.7, fill=CODE_BG)
        add_rounded_box(slide, left, top, col_w, 0.75, fill=banner)
        add_text(slide, left + 0.3, top + 0.1, col_w - 0.6, 0.4, name,
                 size=20, bold=True, color=WHITE)
        add_text(slide, left + 0.3, top + 0.42, col_w - 0.6, 0.35, tag,
                 size=12, color=WHITE)
        add_bullets(slide, left + 0.15, top + 0.95, col_w - 0.3, 3.65,
                    [[(item, {"color": WHITE})] for item in items], size=12)
    set_notes(
        slide,
        """BLOCK / CONFIRM / WARN maps to how you want an AI teammate to behave: never / pause / heads-up.""",
        files_referenced=["standards/SECURITY-GUARDRAILS.md",
                          ".cursor/rules/security.mdc",
                          "templates/CLAUDE.md lines 33–62"],
    )


def build_slide_8_quality(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Code quality standards",
              subtitle="Consistent AI-generated code by default.")

    left_col = [
        [("Before claiming done", {"bold": True, "color": MAGENTA, "no_bullet": True})],
        [("1. Run tests and report results", {"color": WHITE})],
        [("2. Check for lint errors", {"color": WHITE})],
        [("3. Verify build succeeds", {"color": WHITE})],
        [("4. Describe what was tested", {"color": WHITE})],
        [("", {"no_bullet": True})],
        [("Comments", {"bold": True, "color": MAGENTA, "no_bullet": True})],
        [("WHY only — never obvious WHAT", {"color": WHITE})],
        [("No commented-out code", {"color": WHITE})],
        [("Document breaking changes", {"color": WHITE})],
    ]
    right_col = [
        [("Error handling", {"bold": True, "color": MAGENTA, "no_bullet": True})],
        [("Handle all error cases explicitly", {"color": WHITE})],
        [("Meaningful messages with context", {"color": WHITE})],
        [("Never swallow exceptions silently", {"color": WHITE})],
        [("", {"no_bullet": True})],
        [("Structure", {"bold": True, "color": MAGENTA, "no_bullet": True})],
        [("Imports at top; each function does one thing", {"color": WHITE})],
        [("Prefer editing existing files over new", {"color": WHITE})],
        [("Small incremental changes", {"color": WHITE})],
        [("Extensions: Python · TypeScript · template", {"color": WHITE})],
    ]
    add_bullets(slide, 0.7, 2.2, 6.0, 4.2, left_col, size=14, bullet_char="›")
    add_bullets(slide, 7.0, 2.2, 6.0, 4.2, right_col, size=14, bullet_char="›")

    add_rounded_box(slide, 0.7, 6.15, SLIDE_W_IN - 1.4, 0.75,
                    fill=CODE_BG, accent_left=True)
    add_text(slide, 1.0, 6.25, SLIDE_W_IN - 2.0, 0.55,
             "Targets: > 80% test coverage on new code · 0 lint errors · > 90% type coverage · function complexity < 10.",
             size=13, bold=True, color=WHITE)
    set_notes(
        slide,
        """Verification before done is the biggest win. Stops 'I think it works' declarations.""",
        files_referenced=["standards/CODE-QUALITY.md",
                          "templates/CLAUDE.md lines 64–87",
                          "standards/extensions/python.md, typescript.md"],
    )


def build_slide_9_logging(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Structured logging",
              subtitle="Queryable logs. No secrets. Correlation across services.")

    add_text(slide, 0.7, 2.2, 6.0, 0.5, "Why it matters",
             size=18, bold=True, color=MAGENTA)
    items = [
        [("Incident response", {"bold": True, "color": WHITE}),
         ("  — filter by user_id / order_id / correlation_id.", {"color": SUB_TEXT})],
        [("SRE / on-call", {"bold": True, "color": WHITE}),
         ("  — JSON logs feed log aggregators directly.", {"color": SUB_TEXT})],
        [("Compliance", {"bold": True, "color": WHITE}),
         ("  — PII auto-redacted before hitting disk.", {"color": SUB_TEXT})],
        [("Noise reduction", {"bold": True, "color": WHITE}),
         ("  — skip /health and /ready by default.", {"color": SUB_TEXT})],
    ]
    add_bullets(slide, 0.7, 2.75, 6.0, 3.4, items, size=13)

    add_text(slide, 7.0, 2.2, 5.9, 0.5, "Required pattern",
             size=18, bold=True, color=MAGENTA)
    code = (
        '# GOOD — structured, queryable\n'
        'logger.info("order_created",\n'
        '    order_id="ORD-123",\n'
        '    total=149.99)\n'
        '\n'
        '# BAD — string concatenation\n'
        'logger.info(f"Order {id} done")\n'
        '\n'
        '# correlation ID\n'
        'logger = logger.bind(\n'
        '    correlation_id=req_id)'
    )
    add_code_box(slide, 7.0, 2.75, 5.9, 3.4, code)

    add_rounded_box(slide, 0.7, 6.3, SLIDE_W_IN - 1.4, 0.6,
                    fill=CODE_BG, accent_left=True)
    add_text(slide, 1.0, 6.4, SLIDE_W_IN - 2.0, 0.45,
             "Python: structlog  ·  TypeScript: pino  ·  Env: LOG_LEVEL · LOG_FORMAT=json · LOG_SKIP_PATHS · LOG_CORRELATION_HEADER",
             size=12, color=WHITE, font=FONT_MONO)
    set_notes(
        slide,
        """Structured logging pays off during incidents. PII redaction is automatic when using structlog with PIISanitizer.""",
        files_referenced=["standards/LOGGING.md", "docs/LOGGING-GUIDE.md",
                          "README.md lines 119–130"],
    )


def build_slide_10_workflow(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Feature Workflow",
              subtitle="Prevents the #1 AI failure mode: writing code before understanding the problem.")

    # Breadcrumb: Brainstorm -> Spec -> Plan -> Implement (TDD) -> Simplify -> Security Review -> Commit
    breadcrumb_segments = [
        [("Brainstorm", {"color": WHITE, "bold": True}),
         ("  →  ", {"color": SUB_TEXT}),
         ("Spec", {"color": WHITE, "bold": True}),
         ("  →  ", {"color": SUB_TEXT}),
         ("Plan", {"color": WHITE, "bold": True}),
         ("  →  ", {"color": SUB_TEXT}),
         ("Implement (TDD)", {"color": MAGENTA, "bold": True}),
         ("  →  ", {"color": SUB_TEXT}),
         ("Simplify", {"color": WHITE, "bold": True}),
         ("  →  ", {"color": SUB_TEXT}),
         ("Security Review", {"color": WHITE, "bold": True}),
         ("  →  ", {"color": SUB_TEXT}),
         ("Commit", {"color": WHITE, "bold": True})],
    ]
    breadcrumb_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(2.10), Inches(SLIDE_W_IN - 1.4), Inches(0.50)
    )
    tf = breadcrumb_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for segment in breadcrumb_segments[0]:
        text, attrs = segment
        set_run(p.add_run(), text, size=16, font=FONT, **attrs)

    # Left column — Each phase
    add_text(slide, 0.70, 2.85, 6.80, 0.35, "Each phase",
             size=14, bold=True, color=MAGENTA)
    phases = [
        ("Brainstorm", "explore, propose 2–3 approaches."),
        ("Spec", "write validated design to docs/specs/."),
        ("Plan", "bite-sized plan, exact file paths."),
        ("Implement", "TDD: failing test → code → green."),
        ("Simplify", "clarity review, no behavior change."),
        ("Security Review", "scan for 9 patterns."),
        ("Commit", "descriptive message; push."),
    ]
    for i, (name, desc) in enumerate(phases):
        row_top = 3.25 + i * 0.33
        phase_box = slide.shapes.add_textbox(
            Inches(0.70), Inches(row_top), Inches(6.80), Inches(0.30)
        )
        ptf = phase_box.text_frame
        ptf.word_wrap = True
        pp = ptf.paragraphs[0]
        pp.alignment = PP_ALIGN.LEFT
        name_color = MAGENTA if name == "Implement" else WHITE
        set_run(pp.add_run(), name, size=12, bold=True, color=name_color, font=FONT)
        set_run(pp.add_run(), f"  —  {desc}", size=12, color=SUB_TEXT, font=FONT)

    # Right column — Skip rule + How to trigger
    add_text(slide, 7.60, 2.85, 5.10, 0.35, "Skip rule",
             size=14, bold=True, color=MAGENTA)
    add_text(slide, 7.60, 3.25, 5.10, 0.30,
             "Jump straight to Implement for:",
             size=12, color=WHITE)
    skip_items = [
        "Single-file fixes",
        "Typos",
        "Config changes",
        "Changes under 20 lines",
    ]
    for i, item in enumerate(skip_items):
        add_text(slide, 7.80, 3.55 + i * 0.28, 4.90, 0.26,
                 f"•  {item}", size=12, color=SUB_TEXT)

    add_text(slide, 7.60, 4.95, 5.10, 0.35, "How to trigger",
             size=14, bold=True, color=MAGENTA)
    trigger_box = slide.shapes.add_textbox(
        Inches(7.60), Inches(5.30), Inches(5.10), Inches(0.60)
    )
    ttf = trigger_box.text_frame
    ttf.word_wrap = True
    for i, (label, rest) in enumerate([
        ("Claude Code: ", "/feature-dev runs the full 7-phase flow."),
        ("Cursor: ", ".cursor/rules/workflow.mdc enforces it automatically."),
    ]):
        tp = ttf.paragraphs[0] if i == 0 else ttf.add_paragraph()
        tp.alignment = PP_ALIGN.LEFT
        tp.line_spacing = 1.15
        set_run(tp.add_run(), label, size=12, bold=True, color=WHITE, font=FONT)
        set_run(tp.add_run(), rest, size=12, color=SUB_TEXT, font=FONT)

    set_notes(
        slide,
        """Breadcrumb shows the 7-phase flow left-to-right. Implement is highlighted — that's where TDD lives. Left column details each phase; right column covers when to skip ahead and how to trigger the workflow in each IDE.""",
        files_referenced=["templates/CLAUDE.md (Workflow section)", "standards/WORKFLOW.md"],
    )


def build_slide_11_code_review(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Multi-agent /code-review",
              subtitle="Three role-separated reviewers. Uncorrelated contexts. No cross-bias.")
    reviewers = [
        ("Security", MAGENTA, [
            "Hardcoded secrets, tokens, passwords",
            "SQL / NoSQL injection",
            "Broken auth / missing authorization",
            "Exposed errors, stack traces",
            "Weak crypto; eval / exec / shell=True",
            "Sensitive data in logs / responses",
        ], "CRITICAL · HIGH · MEDIUM · LOW"),
        ("Performance", BERRY, [
            "N+1 queries, loops over the DB",
            "Unbounded loops / recursion",
            "Large payloads held in memory",
            "Blocking I/O in async paths",
            "Repeated work that could be cached",
            "Missing pagination",
        ], "HIGH · MEDIUM · LOW"),
        ("Style & Standards", DARK_GRAY, [
            "Functions > 50 lines",
            "Deeply nested logic",
            "Ambiguous names",
            "Missing WHY comments",
            "Dead code / debug statements",
            "Magic numbers; copy-paste dup",
        ], "MEDIUM · LOW"),
    ]
    col_w = 4.05
    top = 2.2
    for i, (title, banner, items, sev) in enumerate(reviewers):
        left = 0.65 + i * (col_w + 0.08)
        add_rounded_box(slide, left, top, col_w, 4.2, fill=CODE_BG)
        add_rounded_box(slide, left, top, col_w, 0.65, fill=banner)
        add_text(slide, left + 0.2, top + 0.13, col_w - 0.4, 0.4, title,
                 size=16, bold=True, color=WHITE)
        add_bullets(slide, left + 0.15, top + 0.8, col_w - 0.3, 2.95,
                    [[(item, {"color": WHITE})] for item in items], size=10)
        add_text(slide, left + 0.2, top + 3.8, col_w - 0.4, 0.3, sev,
                 size=10, color=SUB_TEXT, bold=True)

    add_rounded_box(slide, 0.65, 6.55, SLIDE_W_IN - 1.3, 0.5,
                    fill=CODE_BG, accent_left=True)
    add_text(slide, 0.9, 6.6, SLIDE_W_IN - 1.8, 0.4,
             '"Each subagent sees only the code and its own lens — not the other subagents\' findings."',
             size=12, color=WHITE)
    set_notes(
        slide,
        """ONE half of /code-review — the three parallel reviewers. Next slide covers test coverage + opponent auditor.
Key message: role separation prevents bias.""",
        files_referenced=[".claude/commands/code-review.md Steps 1–3"],
    )


def build_slide_12_testing_auditor(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Test coverage + opponent auditor",
              subtitle="The 'compare' step that makes role-separated review trustworthy.")

    # Left — Test Coverage
    add_rounded_box(slide, 0.65, 2.2, 6.1, 4.2, fill=CODE_BG)
    add_rounded_box(slide, 0.65, 2.2, 6.1, 0.65, fill=MAGENTA)
    add_text(slide, 0.85, 2.33, 5.7, 0.4, "Step 4 — Test Coverage Review",
             size=15, bold=True, color=WHITE)
    add_text(slide, 0.85, 2.93, 5.7, 0.35, "(runs in the main agent)",
             size=11, color=SUB_TEXT)
    test_items = [
        [("Are there tests for changed/new code?", {"color": WHITE})],
        [("Do they cover happy path + edge cases + error paths?", {"color": WHITE})],
        [("Any new public functions or endpoints left untested?", {"color": WHITE})],
        [("Do tests assert meaningful outcomes?", {"color": WHITE})],
        [("Are mocks / stubs used appropriately?", {"color": WHITE})],
    ]
    add_bullets(slide, 0.85, 3.35, 5.7, 1.95, test_items, size=11)
    add_rounded_box(slide, 0.85, 5.35, 5.7, 0.95, fill=BG_NAVY_ALT, accent_left=True)
    add_text(slide, 1.05, 5.45, 5.4, 0.8,
             "If tests are missing, /code-review generates them — one file per changed module, covering happy path + edge cases + all error paths.",
             size=11, bold=True, color=WHITE)

    # Right — Opponent Auditor
    add_rounded_box(slide, 6.85, 2.2, 6.1, 4.2, fill=CODE_BG)
    add_rounded_box(slide, 6.85, 2.2, 6.1, 0.65, fill=BERRY)
    add_text(slide, 7.05, 2.33, 5.7, 0.4, "Step 5 — Opponent Auditor",
             size=15, bold=True, color=WHITE)
    add_text(slide, 7.05, 2.93, 5.7, 0.35,
             "(final subagent — sees findings from A / B / C)",
             size=11, color=SUB_TEXT)
    verdicts = [
        ("✓", "Confirm", "Finding stands.", MAGENTA),
        ("↓", "Downgrade", "Real, but severity too high.", BERRY),
        ("✗", "False positive", "Remove; explain why.", DARK_GRAY),
    ]
    base_top = 3.35
    for i, (mark, title, desc, color) in enumerate(verdicts):
        t = base_top + i * 0.48
        add_text(slide, 7.05, t, 0.4, 0.4, mark,
                 size=16, bold=True, color=color)
        add_text(slide, 7.5, t, 1.5, 0.4, title,
                 size=13, bold=True, color=WHITE)
        add_text(slide, 9.0, t, 3.7, 0.4, desc,
                 size=12, color=SUB_TEXT)
    add_rounded_box(slide, 7.05, 5.0, 5.7, 1.3, fill=BG_NAVY_ALT, accent_left=True)
    add_text(slide, 7.25, 5.1, 5.4, 1.15,
             "Also surfaces anything all three reviewers missed. Final report has an 'Auditor verdict' column on every finding.",
             size=11, color=WHITE)

    # Bottom — why this works
    add_rounded_box(slide, 0.65, 6.55, SLIDE_W_IN - 1.3, 0.5,
                    fill=MAGENTA)
    add_text(slide, 0.9, 6.6, SLIDE_W_IN - 1.8, 0.4,
             "Why this works: role separation keeps each reviewer honest · a compare/auditor step catches misses AND over-flags.",
             size=12, bold=True, color=WHITE)
    set_notes(
        slide,
        """Anchor on two points: (1) Step 4 GENERATES missing tests, not just flags them. (2) Step 5 auditor triages findings so you know what to act on.""",
        files_referenced=[".claude/commands/code-review.md Step 4 (lines 76–87), Step 5 (lines 89–96)"],
    )


def build_slide_13_install_windows(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Install — Windows",
              subtitle="Five steps. No admin rights. ~2 minutes.")
    add_text(slide, 0.7, 2.2, SLIDE_W_IN - 1.4, 0.4,
             "1.  Open Windows PowerShell.", size=16, bold=True, color=WHITE)
    add_text(slide, 0.7, 2.65, SLIDE_W_IN - 1.4, 0.4,
             "2.  Go to your project folder:   cd C:\\path\\to\\your-project",
             size=16, bold=True, color=WHITE)
    add_text(slide, 0.7, 3.1, SLIDE_W_IN - 1.4, 0.4,
             "3.  Paste this command, press Enter:",
             size=16, bold=True, color=WHITE)
    add_code_box(
        slide, 0.7, 3.55, SLIDE_W_IN - 1.4, 1.05,
        "irm https://gitlab.com/tmobile/ere/memory-bank/-/raw/master/scripts/init-memory-bank.ps1 | iex",
    )
    add_text(slide, 0.7, 4.75, SLIDE_W_IN - 1.4, 0.4,
             "4.  Open Cursor or Claude Code in that folder.",
             size=16, bold=True, color=WHITE)
    add_text(slide, 0.7, 5.20, SLIDE_W_IN - 1.4, 0.4,
             "5.  In the AI chat, type exactly:",
             size=16, bold=True, color=WHITE)
    add_rounded_box(slide, 0.7, 5.65, SLIDE_W_IN - 1.4, 0.6,
                    fill=BG_NAVY_ALT, accent_left=True)
    add_text(slide, 1.0, 5.73, SLIDE_W_IN - 2.0, 0.48,
             '"Install memory-bank for this project and follow the setup instructions."',
             size=15, bold=True, color=MAGENTA, font=FONT_MONO)

    add_rounded_box(slide, 0.7, 6.40, SLIDE_W_IN - 1.4, 0.65,
                    fill=CODE_BG, accent_left=True)
    add_text(slide, 1.0, 6.43, SLIDE_W_IN - 2.0, 0.28,
             "You do not need to understand the code. ✓ It worked if the AI reads memory-bank files and asks you to fill them in.",
             size=12, color=WHITE)
    add_text(slide, 1.0, 6.72, SLIDE_W_IN - 2.0, 0.28,
             "If prompted to sign in: use your T-Mobile GitLab login.",
             size=12, color=SUB_TEXT)
    set_notes(
        slide,
        """The install scaffolds into the current working directory — cd into the project first. Fallback if irm/iex is blocked: git clone then run the script locally.""",
        files_referenced=["README.md line 12 (one-liner)",
                          "scripts/init-memory-bank.ps1",
                          "docs/SETUP-GUIDE.md"],
    )


def build_slide_14_install_mac(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Install — Mac", subtitle="Five steps. ~2 minutes.")
    add_text(slide, 0.7, 2.2, SLIDE_W_IN - 1.4, 0.4,
             "1.  Open Terminal.", size=16, bold=True, color=WHITE)
    add_text(slide, 0.7, 2.65, SLIDE_W_IN - 1.4, 0.4,
             "2.  Go to your project folder:   cd /path/to/your-project",
             size=16, bold=True, color=WHITE)
    add_text(slide, 0.7, 3.1, SLIDE_W_IN - 1.4, 0.4,
             "3.  Paste this command, press Return:",
             size=16, bold=True, color=WHITE)
    add_code_box(
        slide, 0.7, 3.55, SLIDE_W_IN - 1.4, 1.05,
        "curl -sSL https://gitlab.com/tmobile/ere/memory-bank/-/raw/master/scripts/init-memory-bank.sh | bash",
    )
    add_text(slide, 0.7, 4.75, SLIDE_W_IN - 1.4, 0.4,
             "4.  Open Cursor or Claude Code in that folder.",
             size=16, bold=True, color=WHITE)
    add_text(slide, 0.7, 5.20, SLIDE_W_IN - 1.4, 0.4,
             "5.  In the AI chat, type exactly:",
             size=16, bold=True, color=WHITE)
    add_rounded_box(slide, 0.7, 5.65, SLIDE_W_IN - 1.4, 0.6,
                    fill=BG_NAVY_ALT, accent_left=True)
    add_text(slide, 1.0, 5.73, SLIDE_W_IN - 2.0, 0.48,
             '"Install memory-bank for this project and follow the setup instructions."',
             size=15, bold=True, color=MAGENTA, font=FONT_MONO)

    add_rounded_box(slide, 0.7, 6.40, SLIDE_W_IN - 1.4, 0.65,
                    fill=CODE_BG, accent_left=True)
    add_text(slide, 1.0, 6.43, SLIDE_W_IN - 2.0, 0.28,
             "You do not need to understand the code. ✓ It worked if the AI reads memory-bank files and asks you to fill them in.",
             size=12, color=WHITE)
    add_text(slide, 1.0, 6.72, SLIDE_W_IN - 2.0, 0.28,
             "If prompted to sign in: use your T-Mobile GitLab login.",
             size=12, color=SUB_TEXT)
    set_notes(
        slide,
        """Same flow as Windows — only the shell differs. curl | bash scaffolds into the current directory.""",
        files_referenced=["README.md line 21", "scripts/init-memory-bank.sh"],
    )


def build_slide_15_first_use(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Recommended first use",
              subtitle="Ten minutes, end-to-end, with an AI teammate.")
    steps = [
        ("1", "Open a project", "In Cursor or Claude Code. Even an existing one."),
        ("2", "Install memory-bank", "Run the one-liner. Fill in the 5 files (~15 min)."),
        ("3", "Ask the AI to read memory-bank", "Quick sanity check."),
        ("4", "Plan a small feature", "Use /feature-dev."),
        ("5", "Let it implement", "TDD: failing test → code → green. Small commits."),
        ("6", "Run /code-review", "Watch reviewers + coverage + auditor."),
        ("7", "Review and commit", "Fix auditor-confirmed findings. Commit."),
    ]
    for i, (num, title, desc) in enumerate(steps):
        top = 2.3 + i * 0.63
        add_rounded_box(slide, 0.7, top, 0.55, 0.5, fill=MAGENTA)
        add_text(slide, 0.7, top + 0.1, 0.55, 0.35, num,
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, 1.4, top + 0.02, 4.4, 0.5, title,
                 size=14, bold=True, color=WHITE)
        add_text(slide, 1.4, top + 0.25, 4.4, 0.35, desc,
                 size=11, color=SUB_TEXT)

    add_rounded_box(slide, 6.1, 2.3, 6.8, 4.4, fill=CODE_BG, accent_top=True)
    add_text(slide, 6.3, 2.42, 6.5, 0.4,
             "What 'good' looks like after the first run",
             size=16, bold=True, color=MAGENTA)
    callouts = [
        [("AI already knows the project", {"bold": True, "color": WHITE}),
         ("  without you re-explaining.", {"color": SUB_TEXT})],
        [("Standards applied consistently", {"bold": True, "color": WHITE}),
         ("  across every file.", {"color": SUB_TEXT})],
        [("Tests exist for new code", {"bold": True, "color": WHITE}),
         ("  — generated by the coverage reviewer if forgotten.", {"color": SUB_TEXT})],
        [("Findings are triaged", {"bold": True, "color": WHITE}),
         ("  — auditor verdict tells you what to fix vs ignore.", {"color": SUB_TEXT})],
        [("Commit is clean", {"bold": True, "color": WHITE}),
         ("  — no secrets, no lint errors, descriptive message.", {"color": SUB_TEXT})],
    ]
    add_bullets(slide, 6.3, 2.9, 6.4, 3.6, callouts, size=12)
    set_notes(
        slide,
        """The first run converts skeptics. Pick a small feature with inputs + outputs + an error path.""",
        files_referenced=["docs/SETUP-GUIDE.md verification section",
                          "training/exercises/01-basic-setup.md"],
    )


def build_slide_16_adoption(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Adoption & next steps",
              subtitle="Start small. Prove value. Expand.")
    phases = [
        ("Week 1", "One repo", "Pick a low-risk project. Run the one-liner. Fill memory-bank."),
        ("Week 2", "One feature", "Use /feature-dev end-to-end. Run /code-review."),
        ("Week 3", "Gather feedback", "What saved time? What was noisy? Adjust."),
        ("Week 4+", "Expand", "3 more repos. Global setup for the team. Add language extensions."),
    ]
    top0 = 2.3
    for i, (when, what, how) in enumerate(phases):
        top = top0 + i * 1.0
        add_rounded_box(slide, 0.7, top, SLIDE_W_IN - 1.4, 0.85,
                        fill=CODE_BG, accent_left=True)
        add_rounded_box(slide, 0.7, top, 1.7, 0.85, fill=MAGENTA)
        add_text(slide, 0.7, top + 0.25, 1.7, 0.4, when,
                 size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, 2.55, top + 0.12, 2.6, 0.4, what,
                 size=16, bold=True, color=WHITE)
        add_text(slide, 2.55, top + 0.45, SLIDE_W_IN - 3.3, 0.4, how,
                 size=12, color=SUB_TEXT)

    add_rounded_box(slide, 0.7, 6.25, SLIDE_W_IN - 1.4, 0.45, fill=MAGENTA)
    add_text(slide, 0.7, 6.30, SLIDE_W_IN - 1.4, 0.35,
             "Safer AI-assisted engineering starts with repeatable defaults.",
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Questions / contact
    add_text(slide, 0.7, 6.78, SLIDE_W_IN - 1.4, 0.28,
             "Questions?  Teams: RE - SkyNet Support - AI Discussion   ·   Email: eric.c.nolan@t-mobile.com",
             size=12, color=SUB_TEXT, align=PP_ALIGN.CENTER)

    set_notes(
        slide,
        """Close with a concrete rollout plan. The closing line is quotable — land it, pause, take questions.
Direct follow-ups to the Teams channel 'RE - SkyNet Support - AI Discussion' or email Eric Nolan (eric.c.nolan@t-mobile.com).""",
        files_referenced=["README.md Contributing + Support", "CONTRIBUTING.md",
                          "training/exercises/"],
    )


def build_slide_handoff(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Handoff: Seamless Across Sessions",
              subtitle="Memory-bank reads automatically. Handoff bridges sessions when context fills up. You never manually re-explain.")

    # Left panel — the seamless loop
    add_rounded_box(slide, 0.7, 2.15, 7.5, 4.5, fill=CODE_BG, accent_top=True)
    add_text(slide, 0.9, 2.28, 7.2, 0.45, "The seamless loop",
             size=18, bold=True, color=MAGENTA)

    loop_lines = [
        ("Session 1", MAGENTA, True),
        ("    AI reads memory-bank/ automatically (no action)", WHITE, False),
        ("    You work. AI tracks decisions, patterns, progress.", WHITE, False),
        ("    Context gets full.", WHITE, False),
        ('    You type: "Handoff"', MAGENTA, True),
        ("    AI writes handoff.md with session state, then stops.", WHITE, False),
        ("", WHITE, False),
        ("Session 2 (fresh chat)", MAGENTA, True),
        ("    AI reads memory-bank/ automatically (still no action)", WHITE, False),
        ("    AI sees handoff.md, reads it, merges into memory-bank.", WHITE, False),
        ("    AI deletes handoff.md.", WHITE, False),
        ("    AI continues exactly where you left off.", WHITE, False),
    ]
    loop_box = slide.shapes.add_textbox(
        Inches(0.9), Inches(2.8), Inches(7.2), Inches(3.75)
    )
    tf = loop_box.text_frame
    tf.word_wrap = True
    for i, (text, color, bold) in enumerate(loop_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.05
        p.space_after = Pt(2)
        r = p.add_run()
        set_run(r, text, size=12, bold=bold, color=color, font=FONT_MONO)

    # Right panel — Threshold by IDE
    add_rounded_box(slide, 8.35, 2.15, 4.25, 4.5, fill=CODE_BG, accent_top=True)
    add_text(slide, 8.55, 2.28, 3.9, 0.45, "Threshold by IDE",
             size=18, bold=True, color=MAGENTA)

    # IDE rows
    rows = [
        ("Claude Code", "65%",
         "Auto-compacts at ~75% — handoff fires before that so nothing is lost."),
        ("Cursor", "80%",
         "Rules re-inject every response; 80% is safe."),
    ]
    top = 2.85
    for name, pct, why in rows:
        add_text(slide, 8.55, top, 2.0, 0.35, name,
                 size=14, bold=True, color=WHITE)
        add_text(slide, 10.6, top, 2.0, 0.35, pct,
                 size=18, bold=True, color=MAGENTA)
        add_text(slide, 8.55, top + 0.4, 3.95, 1.1, why,
                 size=11, color=SUB_TEXT)
        top += 1.65

    # Bottom callout — what you do
    add_rounded_box(slide, 0.7, 6.70, SLIDE_W_IN - 1.4, 0.40, fill=MAGENTA)
    add_text(slide, 0.9, 6.75, SLIDE_W_IN - 1.8, 0.30,
             'What you do:  Type "Handoff" near the context threshold. Memory-bank reading is automatic.',
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    set_notes(
        slide,
        """Key message: memory-bank READING is automatic, always. Handoff is the ONE explicit thing the user types. The merge-and-delete step on Session 2 is what makes it seamless.""",
        files_referenced=["templates/CLAUDE.md lines 175–194 (Handoff Protocol)",
                          "standards/MEMORY-BANK.md Trigger section"],
    )


def build_slide_why_make_it_global(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Why make it global?",
              subtitle="Copy the rules to your user folder once. Every future project inherits them — no per-project setup beyond the memory-bank scaffold.")

    # Left — Why bother?
    add_rounded_box(slide, 0.7, 2.15, 6.0, 4.5, fill=CODE_BG, accent_top=True)
    add_text(slide, 0.9, 2.28, 5.6, 0.45, "Why bother?",
             size=18, bold=True, color=MAGENTA)
    why_items = [
        [("New projects inherit ", {"color": WHITE}),
         ("security, code quality, logging, workflow", {"color": WHITE, "bold": True}),
         (", and Karpathy principles automatically.", {"color": WHITE})],
        [("Slash commands ", {"color": WHITE}),
         ("/feature-dev, /security-review, /code-review", {"color": WHITE, "bold": True, "font": FONT_MONO}),
         (" work in every project.", {"color": WHITE})],
        [("You only scaffold ", {"color": WHITE}),
         ("memory-bank/", {"color": WHITE, "bold": True, "font": FONT_MONO}),
         (" per project — rules are already live.", {"color": WHITE})],
    ]
    add_bullets(slide, 0.9, 2.85, 5.6, 3.7, why_items, size=14)

    # Right — Where rules live
    add_rounded_box(slide, 6.9, 2.15, 6.0, 4.5, fill=CODE_BG, accent_top=True)
    add_text(slide, 7.1, 2.28, 5.6, 0.45, "Where rules live",
             size=18, bold=True, color=MAGENTA)
    locations = (
        "~/.claude/\n"
        "   CLAUDE.md        ← Claude Code rules\n"
        "   AGENTS.md        ← cross-tool rules\n"
        "   commands/*.md    ← slash commands\n"
        "~/.cursor/\n"
        "   rules/*.mdc      ← Cursor rules"
    )
    add_code_box(slide, 7.1, 2.85, 5.6, 3.7, locations)

    # Bottom callout — preview the next two slides
    add_rounded_box(slide, 0.7, 6.85, SLIDE_W_IN - 1.4, 0.25,
                    fill=BG_NAVY_ALT, accent_left=True)
    add_text(slide, 0.9, 6.88, SLIDE_W_IN - 1.8, 0.20,
             "Next two slides: Windows PowerShell, then Mac / Linux Bash. Pick the one for your machine.",
             size=11, color=WHITE)

    set_notes(
        slide,
        """The "why" half of global setup — explains the value (every project inherits security/code-quality/logging/workflow/Karpathy + slash commands) and where the rules live on disk. The next two slides show the actual run-once commands per OS.""",
        files_referenced=["docs/SETUP-GUIDE.md (global setup)",
                          "docs/GLOBAL-RULES-SETUP.md",
                          "docs/CLAUDE-CODE-PLUGINS.md"],
    )


def build_slide_make_it_global(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Make it global (Windows / PowerShell)",
              subtitle="Run once inside the cloned memory-bank repo. Restart Claude Code / Cursor afterwards.")

    add_rounded_box(slide, 0.7, 2.15, SLIDE_W_IN - 1.4, 3.5,
                    fill=CODE_BG, accent_top=True)
    add_text(slide, 0.9, 2.28, SLIDE_W_IN - 1.8, 0.45, "Run once per machine",
             size=18, bold=True, color=MAGENTA)
    win_cmd = (
        'Copy-Item .\\templates\\CLAUDE.md "$env:USERPROFILE\\.claude\\CLAUDE.md"\n'
        'Copy-Item .\\templates\\AGENTS.md "$env:USERPROFILE\\.claude\\AGENTS.md"\n'
        'New-Item -ItemType Directory -Force -Path "$env:USERPROFILE\\.claude\\commands"\n'
        'Copy-Item .\\templates\\claude-commands\\*.md "$env:USERPROFILE\\.claude\\commands\\"\n'
        'New-Item -ItemType Directory -Force -Path "$env:USERPROFILE\\.cursor\\rules"\n'
        'Copy-Item .\\templates\\cursor\\rules\\*.mdc   "$env:USERPROFILE\\.cursor\\rules\\"'
    )
    add_code_box(slide, 0.9, 2.85, SLIDE_W_IN - 1.8, 2.7, win_cmd, size=12)

    add_rounded_box(slide, 0.7, 5.95, SLIDE_W_IN - 1.4, 0.55,
                    fill=BG_NAVY_ALT, accent_left=True)
    add_text(slide, 0.9, 6.02, SLIDE_W_IN - 1.8, 0.45,
             "Restart Claude Code / Cursor to pick up new globals. See docs/GLOBAL-RULES-SETUP.md.",
             size=12, color=WHITE)

    set_notes(
        slide,
        """Windows / PowerShell run-once commands. Full-width code box at default 16pt — long paths fit cleanly. Mac / Linux equivalent on the next slide.""",
        files_referenced=["docs/SETUP-GUIDE.md (global setup)",
                          "docs/GLOBAL-RULES-SETUP.md"],
    )


def build_slide_make_it_global_mac(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Make it global (Mac / Linux)",
              subtitle="Bash equivalent. Same idea — one-time per machine; every future project inherits the rules.")

    add_rounded_box(slide, 0.7, 2.15, SLIDE_W_IN - 1.4, 3.5,
                    fill=CODE_BG, accent_top=True)
    add_text(slide, 0.9, 2.28, SLIDE_W_IN - 1.8, 0.45, "Run once per machine",
             size=18, bold=True, color=MAGENTA)
    mac_cmd = (
        "mkdir -p ~/.claude/commands ~/.cursor/rules\n"
        "cp ./templates/CLAUDE.md            ~/.claude/CLAUDE.md\n"
        "cp ./templates/AGENTS.md            ~/.claude/AGENTS.md\n"
        "cp ./templates/claude-commands/*.md ~/.claude/commands/\n"
        "cp ./templates/cursor/rules/*.mdc   ~/.cursor/rules/"
    )
    add_code_box(slide, 0.9, 2.85, SLIDE_W_IN - 1.8, 2.7, mac_cmd)

    add_rounded_box(slide, 0.7, 5.95, SLIDE_W_IN - 1.4, 0.55,
                    fill=BG_NAVY_ALT, accent_left=True)
    add_text(slide, 0.9, 6.02, SLIDE_W_IN - 1.8, 0.45,
             "Restart Claude Code / Cursor to pick up new globals. See docs/GLOBAL-RULES-SETUP.md.",
             size=12, color=WHITE)

    set_notes(
        slide,
        """Mac/Linux variant of the global setup. Same destination paths, bash idioms instead of PowerShell. Most teams know one or the other — this slide for the bash side.""",
        files_referenced=["docs/SETUP-GUIDE.md (global setup)",
                          "docs/GLOBAL-RULES-SETUP.md"],
    )


def build_slide_quick_commands(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Quick commands",
              subtitle="Daily shortcuts you'll actually use.")

    rows = [
        ("mb status", "Show Memory Bank file sizes and health check"),
        ("mb update", "Update Memory Bank files from the current session"),
        ("mb slim", "Trim activeContext.md to essentials"),
        ("mb commit", "Stage and commit Memory Bank changes"),
        ("Handoff", "Create handoff.md and stop (when context gets full)"),
        ("/feature-dev",
         "Run the full 7-phase workflow (Claude Code)"),
        ("/security-review",
         "Scan the current diff for 9 security patterns (Claude Code)"),
        ("/code-review",
         "Multi-agent review: 3 parallel role subagents + test coverage + opponent auditor"),
        ("/accessibility-review",
         "WCAG 2.1 AA audit of UI code — 9 a11y dimensions, remediation checklist (on-demand)"),
    ]

    # Header row
    col_x = [0.7, 3.7]
    col_w = [2.9, SLIDE_W_IN - 1.4 - 2.9 - 0.1]
    top = 2.15
    row_h = 0.46

    add_rounded_box(slide, col_x[0], top, sum(col_w) + 0.1, row_h - 0.05,
                    fill=MAGENTA)
    add_text(slide, col_x[0] + 0.15, top + 0.08, col_w[0], 0.4, "Command",
             size=14, bold=True, color=WHITE)
    add_text(slide, col_x[1] + 0.1, top + 0.08, col_w[1], 0.4, "What it does",
             size=14, bold=True, color=WHITE)

    for i, (cmd, desc) in enumerate(rows):
        top_r = top + row_h + i * row_h
        if i % 2 == 0:
            add_rounded_box(slide, col_x[0], top_r, sum(col_w) + 0.1,
                            row_h - 0.05, fill=CODE_BG)
        add_text(slide, col_x[0] + 0.15, top_r + 0.1, col_w[0], 0.4, cmd,
                 size=13, bold=True, color=MAGENTA, font=FONT_MONO)
        add_text(slide, col_x[1] + 0.1, top_r + 0.1, col_w[1], 0.4, desc,
                 size=13, color=WHITE)

    set_notes(
        slide,
        """These are the shortcuts for daily use. mb commands are the Memory Bank utility; slash commands run in Claude Code (and the cursor rule equivalents trigger the same workflows).""",
        files_referenced=["docs/QUICK-REFERENCE.md (Quick Commands table)",
                          "scripts/mb.ps1, scripts/mb.sh"],
    )


def build_slide_summary(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Summary", subtitle="What you get from adopting Memory Bank.")

    summary_lines = [
        ("Memory Bank", "Persistent project context"),
        ("Handoff", "Seamless session transitions (65% Claude Code / 80% Cursor)"),
        ("Security Guardrails", "BLOCK / CONFIRM / WARN protection"),
        ("Code Quality", "Consistent AI output, verification before done"),
        ("Structured Logging", "Queryable JSON logs, correlation IDs, PII redaction"),
        ("Feature Workflow",
         "Brainstorm → Spec → Plan → Implement → Simplify → Security → Commit"),
        ("Multi-agent /code-review",
         "Role-separated reviewers + test coverage + opponent auditor"),
        ("Accessibility (UI only)",
         "WCAG 2.1 AA — glob-scoped Cursor rule + /accessibility-review"),
        ("Enterprise hygiene (v1.5)",
         "Prompt injection · Data classification · Ephemeral secrets · OWASP LLM Top 10 · Model governance · Incident runbook"),
        ("Quick Commands", "Easy daily use"),
    ]
    top0 = 2.2
    row_h = 0.46
    for i, (label, desc) in enumerate(summary_lines):
        top = top0 + i * row_h
        add_text(slide, 0.75, top + 0.08, 0.45, 0.4, "✓",
                 size=18, bold=True, color=MAGENTA)
        add_text(slide, 1.25, top + 0.08, 3.7, 0.4, label,
                 size=15, bold=True, color=WHITE)
        add_text(slide, 5.0, top + 0.1, SLIDE_W_IN - 5.4, 0.4, desc,
                 size=13, color=SUB_TEXT)

    set_notes(
        slide,
        """Recap of the eight things the audience now understands. Each line pairs a concept with its outcome. Use this as the 'what did we just cover' slide before Q&A.""",
        files_referenced=["README.md (What's Included)"],
    )


def build_slide_resources(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Resources", subtitle="Where to go next.")

    # Docs list
    docs = [
        ("Setup Guide", "docs/SETUP-GUIDE.md"),
        ("Quick Reference", "docs/QUICK-REFERENCE.md"),
        ("IDE Comparison", "docs/CURSOR-VS-CLAUDE.md"),
        ("Claude Code Plugins + slash commands", "docs/CLAUDE-CODE-PLUGINS.md"),
        ("Core standards",
         "standards/*.md (Memory Bank, Security, Code Quality, Logging, Workflow, Accessibility)"),
        ("Enterprise hygiene (v1.5)",
         "standards/{RULES-FILE-INTEGRITY, DATA-CLASSIFICATION, SECRETS, LLM-TOP-10-MAPPING, MODEL-GOVERNANCE}.md"),
        ("Incident runbook", "templates/INCIDENT-RUNBOOK.md"),
        ("Exercises", "training/exercises/"),
    ]
    top0 = 2.2
    row_h = 0.42
    for i, (label, path) in enumerate(docs):
        top = top0 + i * row_h
        add_text(slide, 0.75, top + 0.04, 4.5, 0.35, label,
                 size=14, bold=True, color=WHITE)
        add_text(slide, 5.3, top + 0.06, SLIDE_W_IN - 5.7, 0.35, path,
                 size=13, color=MAGENTA, font=FONT_MONO)

    # Support box
    support_top = top0 + len(docs) * row_h + 0.3
    add_rounded_box(slide, 0.7, support_top, SLIDE_W_IN - 1.4, 1.15,
                    fill=CODE_BG, accent_left=True)
    add_text(slide, 0.9, support_top + 0.08, 4.0, 0.35, "Support",
             size=16, bold=True, color=MAGENTA)

    add_text(slide, 0.9, support_top + 0.45, 1.2, 0.30, "Teams:",
             size=13, bold=True, color=WHITE)
    add_text(slide, 2.15, support_top + 0.45, SLIDE_W_IN - 2.5, 0.30,
             TEAMS_CHANNEL_NAME, size=13, color=MAGENTA)

    add_text(slide, 0.9, support_top + 0.78, 1.2, 0.30, "GitLab:",
             size=13, bold=True, color=WHITE)
    add_text(slide, 2.15, support_top + 0.78, SLIDE_W_IN - 2.5, 0.30,
             "https://gitlab.com/tmobile/ere/memory-bank",
             size=13, color=MAGENTA, font=FONT_MONO)
    set_notes(
        slide,
        """Final slide — point the audience at the docs and the Teams channel for questions.
Teams channel: RE - SkyNet Support - AI Discussion (direct message Eric Nolan if needed).
GitLab: https://gitlab.com/tmobile/ere/memory-bank""",
        files_referenced=["README.md Support section",
                          "docs/SETUP-GUIDE.md",
                          "docs/QUICK-REFERENCE.md"],
    )


def build_slide_accessibility(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Accessibility (WCAG 2.1 Level AA)",
              subtitle="UI code only. Auto-applies via glob-scoped Cursor rule. On-demand audit via /accessibility-review.")

    # Left panel — always-on rule
    add_rounded_box(slide, 0.7, 2.15, 6.0, 2.3, fill=CODE_BG, accent_top=True)
    add_text(slide, 0.9, 2.28, 5.6, 0.45, "Always-on rule",
             size=16, bold=True, color=MAGENTA)
    add_text(slide, 0.9, 2.75, 5.6, 0.3,
             ".cursor/rules/accessibility.mdc  (glob-scoped)",
             size=11, color=SUB_TEXT, font=FONT_MONO)
    add_code_box(slide, 0.9, 3.1, 5.6, 1.2,
                 ".html  .htm    .jsx   .tsx\n"
                 ".vue   .svelte .astro\n"
                 ".css   .scss   .sass  .less")

    # Right panel — on-demand audit
    add_rounded_box(slide, 6.9, 2.15, 6.0, 2.3, fill=CODE_BG, accent_top=True)
    add_text(slide, 7.1, 2.28, 5.6, 0.45, "On-demand audit",
             size=16, bold=True, color=MAGENTA)
    audit_lines = [
        [("/accessibility-review", {"bold": True, "color": WHITE, "font": FONT_MONO}),
         ("  scans diff, file, or folder.", {"color": SUB_TEXT})],
        [("Emits severity-rated findings (", {"color": SUB_TEXT}),
         ("CRITICAL · HIGH · MEDIUM · LOW", {"color": WHITE, "bold": True}),
         (") × nine dimensions.", {"color": SUB_TEXT})],
        [("Ships a remediation checklist with code-level fixes.", {"color": SUB_TEXT})],
    ]
    add_bullets(slide, 0.9 + 6.2, 2.75, 5.6, 1.55, audit_lines, size=11)

    # Bottom — nine dimensions in a 3x3 grid
    add_text(slide, 0.7, 4.6, SLIDE_W_IN - 1.4, 0.4, "Nine dimensions — all hard requirements",
             size=14, bold=True, color=MAGENTA)

    dims = [
        ("Semantic HTML", "Native elements; heading hierarchy; landmarks."),
        ("ARIA", "Native first; never override; aria-live for dynamic."),
        ("Keyboard", "Reachable; tab order; focus trap; skip-nav."),
        ("Focus indicators", "Custom outline ≥ 3:1 contrast."),
        ("Color & contrast", "4.5:1 normal; 3:1 large; not color-only."),
        ("Forms & inputs", "Programmatic labels; aria-describedby errors."),
        ("Images & media", "alt text; captions; transcripts."),
        ("Motion & animation", "prefers-reduced-motion; no > 3× flashing."),
        ("Testing", "NVDA / JAWS / VoiceOver; accessibility tree."),
    ]
    col_w = 4.05
    row_h = 0.62
    top0 = 5.05
    for i, (label, desc) in enumerate(dims):
        col = i % 3
        row = i // 3
        left = 0.7 + col * (col_w + 0.05)
        top = top0 + row * (row_h + 0.06)
        add_rounded_box(slide, left, top, col_w, row_h, fill=CODE_BG, accent_left=True)
        add_text(slide, left + 0.15, top + 0.06, col_w - 0.3, 0.3, label,
                 size=12, bold=True, color=WHITE)
        add_text(slide, left + 0.15, top + 0.32, col_w - 0.3, 0.3, desc,
                 size=10, color=SUB_TEXT)

    set_notes(
        slide,
        """Accessibility is a conditional standard, not always-on. The Cursor rule auto-activates on UI files only, and the slash command is invoked explicitly. Full standard at standards/ACCESSIBILITY.md. Already propagated to 15 existing team projects alongside Memory-Bank.""",
        files_referenced=["standards/ACCESSIBILITY.md",
                          ".cursor/rules/accessibility.mdc",
                          ".claude/commands/accessibility-review.md"],
    )


def build_slide_enterprise_hygiene(prs):
    slide = add_blank_slide(prs)
    add_title(slide, "Enterprise Hygiene (v1.5)",
              subtitle="Six guardrails that ship with the standard. Each one prevents a real incident — prompt injection, leaked creds, rogue models, or runaway agents.")

    items = [
        ("Rules-file integrity", MAGENTA,
         "Anti-prompt-injection for .cursorrules / CLAUDE.md / AGENTS.md / .mdc / slash-command .md. Glob-scoped rule fires on-open. Rejects invisible Unicode, hidden HTML, guardrail-bypass phrasing.",
         "standards/RULES-FILE-INTEGRITY.md"),
        ("Data Classification", MAGENTA,
         "Four tiers (Public / Internal / Confidential / Restricted) with per-destination rules for prompts, memory-bank, logs, commits, issues. Default upward when unclear.",
         "standards/DATA-CLASSIFICATION.md"),
        ("Ephemeral Secrets", MAGENTA,
         "Vault / AWS SM / Azure KV. No long-lived creds in agent env. Rotate after sessions. Motivated by LiteLLM breach (March 2026).",
         "standards/SECRETS.md"),
        ("OWASP LLM Top 10 mapping", BERRY,
         "All 10 LLM risks mapped to Memory-Bank controls + residual gaps. LLM04/07/08/09 flagged for future passes.",
         "standards/LLM-TOP-10-MAPPING.md"),
        ("Model governance", BERRY,
         "Approved list: Opus 4.7 / Sonnet 4.6 / Haiku 4.5 ✅, GPT requires approval ⚠️, HF public ❌. Version pinning. Canary + review for model upgrades.",
         "standards/MODEL-GOVERNANCE.md"),
        ("Incident runbook + agent controls", BERRY,
         "SEV 1–4 template with AI-involvement checklist. Plus token budgets, loop detection (pause after 3 identical tool calls), 429 handling, MCP monitoring.",
         "templates/INCIDENT-RUNBOOK.md"),
    ]

    # 2x3 grid
    col_w = 6.15
    row_h = 1.55
    top0 = 2.15
    left0 = 0.7
    for i, (label, banner, desc, path) in enumerate(items):
        col = i % 2
        row = i // 2
        left = left0 + col * (col_w + 0.1)
        top = top0 + row * (row_h + 0.1)
        add_rounded_box(slide, left, top, col_w, row_h, fill=CODE_BG, accent_left=True)
        add_rounded_box(slide, left, top, col_w, 0.35, fill=banner)
        add_text(slide, left + 0.2, top + 0.05, col_w - 0.4, 0.3, label,
                 size=13, bold=True, color=WHITE)
        add_text(slide, left + 0.2, top + 0.45, col_w - 0.4, 0.85, desc,
                 size=10, color=SUB_TEXT)
        add_text(slide, left + 0.2, top + row_h - 0.3, col_w - 0.4, 0.25, path,
                 size=9, color=MAGENTA, font=FONT_MONO)

    add_rounded_box(slide, 0.7, 6.5, SLIDE_W_IN - 1.4, 0.5, fill=BG_NAVY_ALT, accent_left=True)
    add_text(slide, 0.9, 6.57, SLIDE_W_IN - 1.8, 0.4,
             "Coverage mapped to OWASP LLM Top 10 (2025). Known gaps tracked in standards/LLM-TOP-10-MAPPING.md — see that file for residual items and planned passes.",
             size=10, color=WHITE)
    set_notes(
        slide,
        """Six additions from the v1.5 enterprise review. Motivated by the 2025 OWASP LLM Top 10, NIST SP 800-218A, CISA/NSA AI guidance, and recent incidents (LiteLLM, PhantomRaven, EchoLeak). User scope was no-bloat, team-success-focused. Residual gaps are explicitly acknowledged in LLM-TOP-10-MAPPING.md rather than silently ignored.""",
        files_referenced=[
            "standards/RULES-FILE-INTEGRITY.md",
            "standards/DATA-CLASSIFICATION.md",
            "standards/SECRETS.md",
            "standards/LLM-TOP-10-MAPPING.md",
            "standards/MODEL-GOVERNANCE.md",
            "templates/INCIDENT-RUNBOOK.md",
            "standards/SECURITY-GUARDRAILS.md (Agent resource controls section)",
        ],
    )


BUILDERS = [
    build_slide_1_title,
    build_slide_2_audience,
    build_slide_3_exec_summary,
    build_slide_4_problem,
    build_slide_5_solution,
    build_slide_6_memory_bank,
    build_slide_handoff,             # NEW — after Memory-Bank
    build_slide_7_security,
    build_slide_8_quality,
    build_slide_accessibility,       # NEW — after Code Quality
    build_slide_enterprise_hygiene,  # NEW — v1.5 hygiene additions
    build_slide_9_logging,
    build_slide_10_workflow,
    build_slide_11_code_review,
    build_slide_12_testing_auditor,
    build_slide_13_install_windows,
    build_slide_14_install_mac,
    build_slide_why_make_it_global,  # NEW — Why bother? + Where rules live
    build_slide_make_it_global,      # Windows-only, full-width code box
    build_slide_make_it_global_mac,  # Mac/Linux-only, full-width code box
    build_slide_15_first_use,
    build_slide_16_adoption,
    build_slide_quick_commands,
    build_slide_summary,
    build_slide_resources,
]


def main():
    global TOTAL_SLIDES
    TOTAL_SLIDES = len(BUILDERS)
    _slide_counter[0] = 0

    prs = Presentation()  # blank; no brand-template inheritance
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    for builder in BUILDERS:
        builder(prs)

    assert len(prs.slides) == len(BUILDERS)
    for i, slide in enumerate(prs.slides, 1):
        notes = slide.notes_slide.notes_text_frame
        text = "\n".join(p.text for p in notes.paragraphs).strip()
        assert text, f"Slide {i} has empty notes"

    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT} — {len(prs.slides)} slides, consistent dark-navy + magenta style.")


if __name__ == "__main__":
    main()
