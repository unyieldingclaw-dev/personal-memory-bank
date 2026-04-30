"""One-off patch script for memory-bank-overview.pptx.

Applies three text edits without regenerating from build_overview_deck.py:
  - Slide 11 (Enterprise Hygiene): rewrite subtitle and bottom callout.
  - Slides 16/17 (Install Windows/Mac): rewrite the bottom-note second line.
  - Slide 22 (Summary): rewrite the Enterprise hygiene summary line.

Output saved as memory-bank-overview-updated.pptx in the same folder.
"""

from pathlib import Path

from pptx import Presentation

SRC = Path(r"C:\Users\ENolan2\cursor\Memory-Bank\memory-bank-overview.pptx")
DST = SRC.with_name("memory-bank-overview-updated.pptx")

# Edits keyed by 0-indexed slide position.
EDITS = {
    10: [  # Slide 11 — Enterprise Hygiene
        (
            "Six additions from the 2025 OWASP / NIST / CISA gap analysis. "
            "No bloat — each one closes a concrete F500 exposure.",
            "Six guardrails that ship with the standard. Each one prevents a "
            "real incident — prompt injection, leaked creds, rogue models, "
            "or runaway agents.",
        ),
        (
            "Residual gaps acknowledged: OWASP LLM04 (Data/Model Poisoning), "
            "LLM07 (System Prompt Leakage), LLM08 (Vector/Embedding), "
            "LLM09 (Misinformation). Tracked in LLM-TOP-10-MAPPING.md — "
            "not fixed this pass.",
            "Coverage mapped to OWASP LLM Top 10 (2025). Known gaps tracked "
            "in standards/LLM-TOP-10-MAPPING.md — see that file for residual "
            "items and planned passes.",
        ),
    ],
    15: [  # Slide 16 — Install Windows
        (
            "You do not need to understand the code. This gives the "
            "assistant the instructions it needs.",
            "You do not need to understand the code. ✓ It worked if "
            "the AI reads memory-bank files and asks you to fill them in.",
        ),
    ],
    16: [  # Slide 17 — Install Mac
        (
            "You do not need to understand the code. This gives the "
            "assistant the instructions it needs.",
            "You do not need to understand the code. ✓ It worked if "
            "the AI reads memory-bank files and asks you to fill them in.",
        ),
    ],
    21: [  # Slide 22 — Summary
        (
            "Rules-file integrity · Data classification · Ephemeral secrets "
            "· OWASP LLM Top 10 · Model governance · Incident runbook",
            "Prompt injection · Data classification · Ephemeral secrets · "
            "OWASP LLM Top 10 · Model governance · Incident runbook",
        ),
    ],
}


def replace_in_frame(tf, old, new):
    """Replace `old` with `new` in a text_frame. Preserves first run's formatting."""
    for para in tf.paragraphs:
        runs = para.runs
        if not runs:
            continue
        full = "".join(r.text for r in runs)
        if old in full:
            runs[0].text = full.replace(old, new)
            for r in runs[1:]:
                r.text = ""
            return True
    return False


def main():
    prs = Presentation(str(SRC))
    report = []
    for idx, pairs in EDITS.items():
        slide = prs.slides[idx]
        for old, new in pairs:
            hit = False
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if replace_in_frame(shape.text_frame, old, new):
                    hit = True
                    break
            report.append((idx + 1, "OK" if hit else "MISS", old[:70].replace("\n", " ")))

    prs.save(str(DST))
    print(f"Saved: {DST}")
    misses = [r for r in report if r[1] == "MISS"]
    for slide_n, status, snippet in report:
        marker = "  " if status == "OK" else "!!"
        print(f"{marker} Slide {slide_n:>2}: {status} | {snippet}...")
    if misses:
        print(f"\n{len(misses)} edit(s) did not match any text on the target slide. "
              f"Verify the source .pptx has the expected 'old' strings.")
        raise SystemExit(1)


if __name__ == "__main__":
    main()
